"""
PPTParser — Parser primário usando python-pptx.
Extrai texto, alternativas e detecta resposta via marcador vermelho.
Versão web: sem pywin32/COM. Usa LibreOffice headless para exportar PDF (opcional).
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
import re
import os
import subprocess
from pathlib import Path

try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False


class PPTParser:
    def __init__(self, ppt_path: str):
        self.ppt_path = ppt_path
        self.prs = Presentation(ppt_path)
        self._pdf_doc = None

    def save(self, output_path: str):
        self.prs.save(output_path)

    def get_shape_rgb(self, color_obj):
        try:
            if not hasattr(color_obj, "type"):
                return None
            if color_obj.type == 1:
                return str(color_obj.rgb)
            return None
        except Exception:
            return None

    def is_target_red(self, color) -> bool:
        rgb_hex = self.get_shape_rgb(color)
        if not rgb_hex or len(rgb_hex) != 6:
            return False
        try:
            r = int(rgb_hex[0:2], 16)
            g = int(rgb_hex[2:4], 16)
            b = int(rgb_hex[4:6], 16)
            return r > 100 and r > g * 1.5 and r > b * 1.5
        except Exception:
            return False

    def get_shape_info(self, shape):
        info = {"is_red": False, "type": "shape", "has_no_fill": False}
        try:
            if shape.shape_type == 13:
                info["type"] = "picture"
                return info
            if shape.shape_type == 6:
                info["type"] = "group"
                return info
            if hasattr(shape, "fill"):
                if shape.fill.type == 0:
                    info["has_no_fill"] = True
                elif shape.fill.type == 1 and hasattr(shape.fill, "fore_color"):
                    if self.is_target_red(shape.fill.fore_color):
                        info["is_red"] = True
            if hasattr(shape, "line"):
                try:
                    if self.is_target_red(shape.line.color):
                        info["is_red"] = True
                except Exception:
                    pass
        except Exception:
            pass
        return info

    def iter_all_shapes(self, shapes):
        for shape in shapes:
            if shape.shape_type == 6:
                for sub in self.iter_all_shapes(shape.shapes):
                    yield sub
            else:
                yield shape

    def _get_pdf_doc(self):
        """Converte PPTX para PDF via LibreOffice headless (sem pywin32)."""
        if self._pdf_doc:
            return self._pdf_doc
        if not FITZ_AVAILABLE:
            return None
        try:
            abs_ppt = str(Path(self.ppt_path).resolve())
            output_dir = str(Path(abs_ppt).parent)
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", output_dir, abs_ppt],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0:
                pdf_path = abs_ppt.replace(".pptx", ".pdf")
                if Path(pdf_path).exists():
                    self._pdf_doc = fitz.open(pdf_path)
                    return self._pdf_doc
        except Exception as e:
            print(f"[PPTParser] PDF conversion failed: {e}")
        return None

    def detect_correct_answer_visual(self, slide, alt_block_candidates=None):
        """Detecta resposta correta via marcador vermelho (coordenadas)."""
        try:
            all_shapes_flat = list(self.iter_all_shapes(slide.shapes))

            alt_block = None
            if alt_block_candidates:
                alt_block = max(alt_block_candidates, key=lambda s: len(s.text_frame.paragraphs))
            else:
                max_paras = 0
                for shape in all_shapes_flat:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        if len(shape.text_frame.paragraphs) > max_paras:
                            max_paras = len(shape.text_frame.paragraphs)
                            alt_block = shape

            if not alt_block:
                return None

            red_markers = []
            for shape in all_shapes_flat:
                if shape == alt_block:
                    continue
                if shape.height > 6000000:
                    continue
                info = self.get_shape_info(shape)
                if info["is_red"] and info["type"] != "picture":
                    red_markers.append(shape)

            if not red_markers:
                return None

            # Tentativa 0: Leitura visual via PDF (requer LibreOffice)
            pdf_doc = self._get_pdf_doc()
            if pdf_doc and FITZ_AVAILABLE:
                try:
                    slide_idx = next(
                        (i for i, s in enumerate(self.prs.slides) if s == slide), -1
                    )
                    if 0 <= slide_idx < len(pdf_doc):
                        page = pdf_doc[slide_idx]
                        for marker in red_markers:
                            rect = fitz.Rect(
                                marker.left / 12700.0,
                                marker.top / 12700.0,
                                (marker.left + marker.width) / 12700.0,
                                (marker.top + marker.height) / 12700.0,
                            )
                            vis_text = page.get_textbox(rect).strip().upper()
                            m = re.search(r"([A-E])", vis_text)
                            if m:
                                return m.group(1)
                except Exception:
                    pass

            # Tentativa 1: caixas separadas para cada alternativa
            letter_boxes = []
            for shape in all_shapes_flat:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip().upper()
                    if re.match(r"^([A-E])\s*[).\-]", txt):
                        letter_boxes.append(shape)

            if len(letter_boxes) >= 3:
                letter_boxes.sort(key=lambda s: s.top)
                best_letter = None
                min_score = float("inf")
                for marker in red_markers:
                    marker_cy = marker.top + marker.height / 2
                    for shape in letter_boxes:
                        s_cy = shape.top + shape.height / 2
                        score = abs(marker_cy - s_cy)
                        if score < min_score:
                            min_score = score
                            m = re.search(r"^([A-E])", shape.text_frame.text.strip().upper())
                            if m:
                                best_letter = m.group(1)
                if best_letter:
                    return best_letter

            # Tentativa 2: Bloco único com gradação de posição
            alt_y = alt_block.top
            alt_h = alt_block.height
            if alt_h == 0:
                return None

            raw_paras = [p.text for p in alt_block.text_frame.paragraphs]
            para_lines = []
            alt_map = {}

            text_indices = [i for i, t in enumerate(raw_paras) if t.strip()]
            alt_indices = text_indices[-5:] if len(text_indices) >= 5 else text_indices

            for i, txt in enumerate(raw_paras):
                m = re.search(r"^\s*([A-E])[\)\.\-]", txt)
                if m:
                    alt_map[i] = m.group(1).upper()
                if not txt.strip():
                    para_lines.append(1.5)
                elif i in alt_indices:
                    para_lines.append(1.0 + len(txt) / 90.0)
                else:
                    lines = sum(max(1.0, len(seg) // 85 + 1.0) for seg in txt.split("\n"))
                    para_lines.append(lines)

            if sum(para_lines) == 0:
                return None

            unit_h = alt_h / sum(para_lines)
            current_y = alt_y
            para_coords: dict[int, float] = {}
            for i, lines in enumerate(para_lines):
                ph = lines * unit_h
                para_coords[i] = current_y + ph / 2.0
                current_y += ph

            for marker in red_markers:
                marker_cy = marker.top + marker.height / 2
                best_letter = None
                min_dist = float("inf")

                if len(alt_map) >= 3:
                    for idx, letter in alt_map.items():
                        dist = abs(marker_cy - para_coords[idx])
                        if dist < min_dist:
                            min_dist = dist
                            best_letter = letter
                else:
                    letters = ["A", "B", "C", "D", "E"]
                    for i, p_idx in enumerate(alt_indices):
                        dist = abs(marker_cy - para_coords[p_idx])
                        if dist < min_dist:
                            min_dist = dist
                            best_letter = letters[i] if i < len(letters) else None

                if best_letter:
                    return best_letter

        except Exception as e:
            print(f"[PPTParser] detect error: {e}")
        return None

    def get_slide_data(self, gabarito_path: str | None = None) -> list[dict]:
        """Extrai questões e alternativas de todos os slides."""
        gabarito: dict | list | None = None
        if gabarito_path and Path(gabarito_path).exists():
            gabarito = self._load_gabarito(gabarito_path)

        slides_data = []
        buffered_question: dict | None = None

        for slide_idx, slide in enumerate(self.prs.slides):
            all_shapes = list(self.iter_all_shapes(slide.shapes))
            text_blocks = [
                s for s in all_shapes
                if s.has_text_frame and s.text_frame.text.strip()
            ]

            question_text = ""
            alternatives: list[str] = []
            alt_block_candidates: list = []

            for shape in text_blocks:
                text = shape.text_frame.text.strip()
                lines = [p.text.strip() for p in shape.text_frame.paragraphs]

                if any(re.match(r"^[A-E]\s*[\)\.]\s*.+", ln) for ln in lines if ln):
                    alt_block_candidates.append(shape)
                    for ln in lines:
                        if ln:
                            alternatives.append(ln)
                else:
                    if len(text) > 20:
                        question_text += (" " + text if question_text else text)

            # Se não encontrou alternativas com prefixo A/B/C, tenta extrair do texto
            if not alternatives and question_text:
                extracted_q, extracted_alts = self._extract_alts_from_text(question_text)
                if extracted_alts:
                    question_text = extracted_q
                    alternatives = extracted_alts

            # Tentar detectar resposta
            correct_answer = None
            if gabarito:
                correct_answer = self._get_gabarito_answer(gabarito, slide_idx, len(slides_data))

            if correct_answer is None and alt_block_candidates:
                correct_answer = self.detect_correct_answer_visual(slide, alt_block_candidates)
            elif correct_answer is None:
                correct_answer = self.detect_correct_answer_visual(slide)

            # Detectar número da questão
            q_num = None
            for shape in text_blocks:
                text = shape.text_frame.text.strip()
                m = re.search(r"(?:quest[aã]o|q\.?)\s*(\d+)", text, re.IGNORECASE)
                if not m:
                    m = re.match(r"^(\d+)\s*[\.)]", text)
                if m:
                    q_num = int(m.group(1))
                    break

            if not alternatives and buffered_question:
                # Se o slide atual tem texto substancial (> 200 chars), é uma nova questão
                # Não é continuação — flush o buffer e trata como nova questão
                if len(question_text) > 200:
                    slides_data.append(buffered_question)
                    buffered_question = None
                    # Cai para criação de nova questão abaixo
                else:
                    buffered_question["question"] += "\n" + question_text
                    continue

            if buffered_question and alternatives:
                buffered_question["alternatives"] = alternatives
                buffered_question["correct_answer"] = correct_answer
                slides_data.append(buffered_question)
                buffered_question = None
                continue

            if question_text and not alternatives:
                buffered_question = {
                    "slide_index": slide_idx,
                    "question_number": q_num or (len(slides_data) + 1),
                    "question": question_text,
                    "alternatives": [],
                    "correct_answer": correct_answer,
                }
                continue

            if question_text or alternatives:
                slides_data.append({
                    "slide_index": slide_idx,
                    "question_number": q_num or (len(slides_data) + 1),
                    "question": question_text,
                    "alternatives": alternatives,
                    "correct_answer": correct_answer,
                })

        if buffered_question:
            slides_data.append(buffered_question)

        return slides_data

    def _extract_alts_from_text(self, text: str) -> tuple[str, list[str]]:
        """
        Tenta extrair alternativas de texto quando não há prefixo A)/B)/C).
        Heurísticas:
        1. Texto termina com 'assinale...', 'qual é...', 'qual das...' etc → linhas seguintes são alternativas
        2. Últimas 4-5 linhas são curtas e parecidas em tamanho → provavelmente alternativas
        """
        lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
        if len(lines) < 4:
            return text, []

        # Heurística 1: detectar onde a questão termina e as alternativas começam
        question_end_patterns = [
            r"assinale\s+a\s+alternativa",
            r"assinale\s+a\s+opção",
            r"qual\s+[eé]\s+a\s+",
            r"qual\s+das\s+",
            r"indique\s+a\s+",
            r"marque\s+a\s+",
            r"escolha\s+a\s+",
            r"identifique\s+",
            r"\?\s*$",
        ]

        split_idx = -1
        for i, line in enumerate(lines):
            for pattern in question_end_patterns:
                if re.search(pattern, line, re.IGNORECASE):
                    split_idx = i + 1
                    break
            if split_idx > 0:
                break

        if 0 < split_idx < len(lines):
            alt_lines = lines[split_idx:]
            # Valida: 2 a 8 alternativas, nenhuma muito longa
            if 2 <= len(alt_lines) <= 8 and all(len(a) < 300 for a in alt_lines):
                return "\n".join(lines[:split_idx]), alt_lines

        # Heurística 2: últimas 4 linhas se forem curtas (< 200 chars cada)
        if len(lines) >= 6:
            potential_alts = lines[-4:]
            if all(len(a) < 200 for a in potential_alts):
                return "\n".join(lines[:-4]), potential_alts

        return text, []

    def _load_gabarito(self, path: str) -> dict | list | None:
        try:
            with open(path, encoding="utf-8") as f:
                content = f.read().strip()
            result: dict[int, str] = {}
            seq: list[str] = []
            for line in content.splitlines():
                line = line.strip()
                if not line:
                    continue
                m = re.match(r"^\s*(?:Q\.?\s*)?(\d+)[\.\):\s]+([A-Ea-e])", line)
                if m:
                    result[int(m.group(1))] = m.group(2).upper()
                else:
                    m2 = re.match(r"^([A-Ea-e])\s*$", line)
                    if m2:
                        seq.append(m2.group(1).upper())
            if result:
                return result
            if seq:
                return seq
        except Exception:
            pass
        return None

    def _get_gabarito_answer(self, gabarito, slide_idx: int, q_count: int) -> str | None:
        if isinstance(gabarito, dict):
            key = q_count + 1
            return gabarito.get(key)
        if isinstance(gabarito, list):
            if q_count < len(gabarito):
                return gabarito[q_count]
        return None
