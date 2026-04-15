"""
PPTX Extractor Service — porta a lógica do emrautomacao (Excel Add-in) para Python.

Fluxo:
  1. Abre o .pptx com python-pptx
  2. Lê formas nomeadas ("codigo", "enunciado", "alternativas") de cada slide
  3. Lê as notas de cada slide (dica, resumo, justificativas, resposta, especialidade, assunto)
  4. Detecta inverter_comentario via regex
  5. Atribui professores por faixa de questões
  6. Retorna lista de dicts com as 27 colunas prontas para CSV
"""
from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Pt


# ── Constantes ────────────────────────────────────────────────────────────────

COLUNAS = [
    "banca", "ano", "tipo_de_prova", "codigo", "enunciado",
    "alternativa_correta", "anulada",
    "alternativa_a", "alternativa_b", "alternativa_c", "alternativa_d", "alternativa_e",
    "imagem", "codigo_vimeo", "dica_emr", "inverter_comentario", "descricao_comentario",
    "justificativa_alternativa_a", "justificativa_alternativa_b",
    "justificativa_alternativa_c", "justificativa_alternativa_d",
    "justificativa_alternativa_e",
    "nome_professor_comentario", "nome_professor_video", "grande_area_oficial",
    "especialidade", "assunto",
]

PROFESSORES = [
    "Ana Beatriz", "Amadeu Marinho", "Amanda Truta", "André Lafayette",
    "Arthur Régis", "Bruna Holanda", "Caio Atanasio", "Caio Pares",
    "Camila Firme", "Catarina Braga", "Carlos Piscoya", "Dahra Teles",
    "Débora Leite", "Danilo de Holanda", "Filipe Marinho", "Francilberto Souza",
    "Gabriel Melo", "Gabrielle Souza", "George Carvalho", "Henrique Nascimento",
    "Igor Melo", "Jéssica Cesario", "João Manoel", "Juliana Amorim",
    "Lucas Reis", "Mariana Muniz", "Marcos Paixão", "Mariana Silva",
    "Nathalia Almeida", "Matheus Brandt", "Mylena Medeiros", "Rafael Viana",
    "Raphael Burlamaqui", "Raphaella Leite", "Rebecca Castelo Branco", "Rhaissa Holanda",
    "Rodrigo Suassuna", "Suzana Leite", "Thais Rodrigues", "Thiago Arruda",
    "Thiago Oliveira", "Thiago Silva", "Valeria Lafayette",
]

GRANDES_AREAS = [
    "Cirurgia",
    "Clínica Médica",
    "Geral",
    "Ginecologia e Obstetrícia",
    "Medicina Preventiva",
    "Pediatria",
]


# ── Dataclasses ───────────────────────────────────────────────────────────────

@dataclass
class ProfessorBloco:
    range_nums: Optional[Dict[int, bool]]   # None = todas as questões
    nome_comentario: str
    nome_video: str


@dataclass
class NotasSlide:
    dica: str = ""
    resumo: str = ""
    just: Dict[str, str] = field(default_factory=lambda: {k: "" for k in "ABCDE"})
    resposta: str = ""
    especialidade: str = ""
    assunto: str = ""
    anulada_por_nota: bool = False


# ── Utilidades ────────────────────────────────────────────────────────────────

def to_slug(s: str) -> str:
    """Converte string para slug ASCII lowercase (igual ao toSlug() do JS)."""
    s = s.lower().strip()
    s = unicodedata.normalize("NFD", s)
    s = "".join(c for c in s if unicodedata.category(c) != "Mn")
    s = re.sub(r"\s+", "-", s)
    s = re.sub(r"[^a-z0-9-]", "", s)
    return s


def parse_range(s: str) -> Optional[Dict[int, bool]]:
    """Converte string de range '1-10, 15, 20-25' em dict {num: True}."""
    if not s or not s.strip():
        return None
    nums: Dict[int, bool] = {}
    for parte in s.split(","):
        parte = parte.strip()
        m = re.match(r"^(\d+)\s*-\s*(\d+)$", parte)
        if m:
            for n in range(int(m.group(1)), int(m.group(2)) + 1):
                nums[n] = True
        elif re.match(r"^\d+$", parte):
            nums[int(parte)] = True
    return nums or None


def _limpar(s: str) -> str:
    return re.sub(r"^[\s:]+", "", (s or "")).strip()


# ── Extração de formas ────────────────────────────────────────────────────────

def _get_shape_text(slide, name: str) -> str:
    """Retorna o texto de uma forma nomeada do slide."""
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            return shape.text_frame.text.strip()
    return ""


def _get_alternativas(slide) -> List[str]:
    """Retorna lista de alternativas da forma 'alternativas' (uma por parágrafo)."""
    for shape in slide.shapes:
        if shape.name == "alternativas" and shape.has_text_frame:
            result = []
            for para in shape.text_frame.paragraphs:
                linha = para.text.strip()
                if linha:
                    result.append(linha)
            return result
    return []


def _slide_has_anulada(slide) -> bool:
    """Verifica se qualquer forma do slide contém 'ANULADA'."""
    for shape in slide.shapes:
        if shape.has_text_frame and "ANULADA" in shape.text_frame.text.upper():
            return True
    return False


# ── Parsing de notas ──────────────────────────────────────────────────────────

def _buscar_setor(texto: str, inicio_re: re.Pattern, fins_re: re.Pattern) -> str:
    """Extrai seção de texto entre dois marcadores regex (porta do buscarSetor JS)."""
    mi = inicio_re.search(texto)
    if not mi:
        return ""
    resto = texto[mi.end():]
    mf = fins_re.search(resto)
    return (resto[:mf.start()] if mf else resto).strip()


def _parse_notas(slide) -> NotasSlide:
    """Extrai e estrutura o conteúdo das notas de um slide."""
    result = NotasSlide()

    if not slide.has_notes_slide:
        return result

    notes_tf = slide.notes_slide.notes_text_frame
    if not notes_tf:
        return result

    full_text = " ".join(
        para.text for para in notes_tf.paragraphs
    ).replace("\n", " ")
    full_text = re.sub(r"\s+", " ", full_text).strip()

    if not full_text:
        return result

    re_dica   = re.compile(r"Dica\s+de\s+Prova:?", re.I)
    re_resumo = re.compile(r"Resumo\s+do\s+Tema:?", re.I)
    re_just_a = re.compile(r"Coment[aá]rio\s+Alternativa\s+por\s+Alternativa:?", re.I)
    re_just_b = re.compile(r"Coment[aá]rio\s*:", re.I)
    re_resp   = re.compile(r"Resposta\s+correta:?\s*Letra", re.I)
    re_esp    = re.compile(r"Especialidade:?", re.I)
    re_ass    = re.compile(r"Assunto:?", re.I)

    re_just = re_just_a if re_just_a.search(full_text) else re_just_b

    fim_dica   = re.compile(re_resumo.pattern + "|" + re_just.pattern, re.I)
    fim_resumo = re.compile(re_just.pattern, re.I)

    result.dica   = _buscar_setor(full_text, re_dica, fim_dica)
    result.resumo = _buscar_setor(full_text, re_resumo, fim_resumo)

    bloco_just = _buscar_setor(full_text, re_just, re_resp)
    for letra in "ABCDE":
        re_letra = re.compile(
            r"Letra\s+" + letra + r"[:\s\-]+(.+?)(?=Letra\s+[A-E][:\s\-]|Resposta\s+correta|$)",
            re.I | re.DOTALL,
        )
        mj = re_letra.search(bloco_just)
        if mj:
            just_txt = mj.group(1).strip()
            just_txt = re.sub(r"^Letra\s+[A-E][:\s\-]+", "", just_txt, flags=re.I).strip()
            result.just[letra] = just_txt

    m_resp = re.search(r"Resposta\s+correta\s*:?\s*Letra\s+([A-E])", full_text, re.I)
    if m_resp:
        result.resposta = m_resp.group(1).upper()
    elif re.search(r"Resposta\s+correta\s*:?[^A-E]*anulada", full_text, re.I):
        result.anulada_por_nota = True

    fim_esp = re.compile(re_ass.pattern + "|" + re_dica.pattern, re.I)
    result.especialidade = _buscar_setor(full_text, re_esp, fim_esp)
    result.assunto       = _buscar_setor(full_text, re_ass, re_dica)

    if re.search(r"anulada", full_text, re.I):
        result.anulada_por_nota = True

    return result


# ── Inverter comentário ───────────────────────────────────────────────────────

def _inverter_comentario(enunciado: str) -> str:
    """Detecta se a questão pede a alternativa INCORRETA/FALSA."""
    txt = enunciado or ""
    partes = re.split(r"[.?!]\s+", txt)
    comando = partes[-1] if partes else txt

    # Guard: questão V/F não inverte
    if re.search(
        r"verdadeir.{0,30}fals|fals.{0,30}verdadeir"
        r"|assinalando.{0,5}v.{0,5}f|marque.{0,5}v.{0,5}f",
        comando, re.I
    ):
        return "F"

    padroes = [
        re.compile(r"incorreta", re.I),
        re.compile(r"incorretas", re.I),
        re.compile(r"incorreto", re.I),
        re.compile(r"incorretos", re.I),
        re.compile(r"falsa", re.I),
        re.compile(r"falsas", re.I),
        re.compile(r"errada", re.I),
        re.compile(r"erradas", re.I),
        re.compile(r"n[aã]o \w+ corret", re.I),
        re.compile(r"n[aã]o \w+ verdadeir", re.I),
        re.compile(r"n[aã]o \w+ adequad", re.I),
        re.compile(r"assinale .*n[aã]o.{0,20}(corret|verdadeir|adequad)", re.I),
        re.compile(r"qual .{0,40}n[aã]o (faz|é|est|corres|pertence|integra|compõe|comp)", re.I),
        re.compile(r"qual .{0,15}n[aã]o.{0,20}(corret|verdadeir|adequad)", re.I),
        re.compile(r"marque .*n[aã]o.{0,20}(corret|verdadeir|adequad)", re.I),
    ]
    for p in padroes:
        if p.search(comando):
            return "V"
    return "F"


# ── Metadados ─────────────────────────────────────────────────────────────────

def _extrair_metadados(prs: Presentation) -> Dict[str, str]:
    """
    Extrai banca, ano e tipo da capa (primeiros 3 slides).
    Estratégia:
      1. Tenta padrão "SIGLA 2026 Tipo" em uma única linha/parágrafo.
      2. Fallback: coleta todo o texto do slide e procura o primeiro ano (20xx/19xx).
    """
    meta = {"banca": "", "ano": "", "tipo": ""}

    for slide in list(prs.slides)[:3]:
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt:
                    continue
                slide_texts.append(txt)
                # Estratégia 1: "SIGLA 2026 Tipo de Prova" na mesma linha
                m = re.match(r"^([A-ZÁÉÍÓÚ\-]{2,})\s+(\d{4})\s+(.+)$", txt)
                if m and not meta["ano"]:
                    meta["banca"] = m.group(1)
                    meta["ano"]   = m.group(2)
                    meta["tipo"]  = m.group(3).strip()

        # Estratégia 2: ano em qualquer texto do slide (20xx ou 19xx)
        if not meta["ano"] and slide_texts:
            full = " ".join(slide_texts)
            m_ano = re.search(r"\b(20\d{2}|19\d{2})\b", full)
            if m_ano:
                meta["ano"] = m_ano.group(1)

        if meta["ano"]:
            break

    return meta


# ── Professores ───────────────────────────────────────────────────────────────

def _resolver_professores(q_idx_1based: int, blocos: List[ProfessorBloco]) -> Tuple[str, str]:
    """Retorna (nome_comentario, nome_video) para a questão de número q_idx_1based."""
    com, vid = "", ""
    for bloco in blocos:
        no_range = bloco.range_nums is None
        in_range = bloco.range_nums is not None and bloco.range_nums.get(q_idx_1based)
        if no_range or in_range:
            if not com and bloco.nome_comentario:
                com = bloco.nome_comentario
            if not vid and bloco.nome_video:
                vid = bloco.nome_video
    return com, vid


# ── Extração principal ────────────────────────────────────────────────────────

def extrair_questoes_pptx(
    pptx_path: str,
    grande_area: str = "",
    professores: Optional[List[ProfessorBloco]] = None,
    filtro_exportacao: Optional[str] = None,
    logger: Callable = print,
) -> List[Dict]:
    """
    Lê o .pptx e retorna lista de dicts com as 27 colunas do Manager.

    Args:
        pptx_path: caminho do arquivo .pptx
        grande_area: ex. "Cirurgia", "Clínica Médica"
        professores: lista de ProfessorBloco com ranges e nomes
        filtro_exportacao: string de range ex. "1-10, 15" (None = todas)
        logger: função de log
    """
    if not Path(pptx_path).exists():
        raise FileNotFoundError(f"Arquivo não encontrado: {pptx_path}")

    prs = Presentation(pptx_path)
    meta = _extrair_metadados(prs)
    blocos = professores or []
    filtro = parse_range(filtro_exportacao) if filtro_exportacao else None

    logger(f"Slides totais: {len(prs.slides)} | Banca detectada: {meta['banca']} {meta['ano']}")

    questoes_map: Dict[str, dict] = {}
    ordem: List[str] = []

    # Pula os 2 primeiros slides (capa/intro)
    slides = list(prs.slides)[2:]

    for s_idx, slide in enumerate(slides):
        codigo = _get_shape_text(slide, "codigo")
        if not codigo:
            continue

        notas = _parse_notas(slide)
        enunciado_raw = _get_shape_text(slide, "enunciado")
        # Remove prefixo "(N) - " se existir
        enunciado = re.sub(r"^\([^)]+\)\s*-\s*", "", enunciado_raw).strip()
        alts = _get_alternativas(slide)
        anulada = _slide_has_anulada(slide) or notas.anulada_por_nota

        if codigo not in questoes_map:
            questoes_map[codigo] = {
                "codigo": codigo,
                "enunciado": enunciado,
                "alts": alts,
                "gabarito": notas.resposta,
                "dica": notas.dica,
                "resumo": notas.resumo,
                "just": dict(notas.just),
                "anulada": anulada,
                "especialidade": notas.especialidade,
                "assunto": notas.assunto,
            }
            ordem.append(codigo)
        else:
            q = questoes_map[codigo]
            if enunciado and len(enunciado) > 10:
                q["enunciado"] += " " + enunciado
            if alts and not q["alts"]:
                q["alts"] = alts
            if notas.resposta:
                q["gabarito"] = notas.resposta
            if anulada:
                q["anulada"] = True
            if notas.dica and not q["dica"]:
                q["dica"] = notas.dica
            if notas.resumo and not q["resumo"]:
                q["resumo"] = notas.resumo
            for lt in "ABCDE":
                if notas.just.get(lt) and not q["just"].get(lt):
                    q["just"][lt] = notas.just[lt]
            if notas.especialidade and not q["especialidade"]:
                q["especialidade"] = notas.especialidade
            if notas.assunto and not q["assunto"]:
                q["assunto"] = notas.assunto

    logger(f"Questões únicas encontradas: {len(ordem)}")

    resultado: List[Dict] = []
    for idx, codigo in enumerate(ordem):
        q = questoes_map[codigo]
        num_questao = idx + 1

        # Aplica filtro de exportação
        if filtro and not filtro.get(num_questao):
            continue

        com, vid = _resolver_professores(num_questao, blocos)

        row = {
            "banca":                       "",
            "ano":                         meta["ano"],
            "tipo_de_prova":               "",
            "codigo":                      q["codigo"],
            "enunciado":                   q["enunciado"],
            "alternativa_correta":         q["gabarito"],
            "anulada":                     "V" if q["anulada"] else "F",
            "alternativa_a":               q["alts"][0] if len(q["alts"]) > 0 else "",
            "alternativa_b":               q["alts"][1] if len(q["alts"]) > 1 else "",
            "alternativa_c":               q["alts"][2] if len(q["alts"]) > 2 else "",
            "alternativa_d":               q["alts"][3] if len(q["alts"]) > 3 else "",
            "alternativa_e":               q["alts"][4] if len(q["alts"]) > 4 else "",
            "imagem":                      "",
            "codigo_vimeo":                "",
            "dica_emr":                    _limpar(q["dica"]),
            "inverter_comentario":         _inverter_comentario(q["enunciado"]),
            "descricao_comentario":        _limpar(q["resumo"]),
            "justificativa_alternativa_a": _limpar(q["just"].get("A", "")),
            "justificativa_alternativa_b": _limpar(q["just"].get("B", "")),
            "justificativa_alternativa_c": _limpar(q["just"].get("C", "")),
            "justificativa_alternativa_d": _limpar(q["just"].get("D", "")),
            "justificativa_alternativa_e": _limpar(q["just"].get("E", "")),
            "nome_professor_comentario":   com,
            "nome_professor_video":        vid,
            "grande_area_oficial":         (grande_area or "").strip(),
            "especialidade":               _limpar(q["especialidade"]),
            "assunto":                     _limpar(q["assunto"]),
        }
        resultado.append(row)

    logger(f"Total exportado: {len(resultado)} questões.")
    return resultado


# ── Geração de Excel ──────────────────────────────────────────────────────────

def questoes_to_xlsx_bytes(questoes: List[Dict]) -> bytes:
    """
    Converte lista de questões (dicts com 27 colunas) para bytes de um .xlsx
    pronto para download, idêntico ao formato do emrautomacao.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from io import BytesIO

    wb = Workbook()
    ws = wb.active
    ws.title = "Questões"

    header_font  = Font(bold=True, color="FFFFFF")
    header_fill  = PatternFill(fill_type="solid", fgColor="2E75B6")
    wrap_align   = Alignment(wrap_text=True, vertical="top")
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # Cabeçalho
    ws.append(COLUNAS)
    for cell in ws[1]:
        cell.font      = header_font
        cell.fill      = header_fill
        cell.alignment = header_align

    ws.row_dimensions[1].height = 30

    # Larguras razoáveis por coluna
    col_widths = {
        "banca": 10, "ano": 6, "tipo_de_prova": 14, "codigo": 22,
        "enunciado": 60, "alternativa_correta": 8, "anulada": 7,
        "alternativa_a": 40, "alternativa_b": 40, "alternativa_c": 40,
        "alternativa_d": 40, "alternativa_e": 40,
        "imagem": 12, "codigo_vimeo": 14,
        "dica_emr": 40, "inverter_comentario": 8, "descricao_comentario": 60,
        "justificativa_alternativa_a": 50, "justificativa_alternativa_b": 50,
        "justificativa_alternativa_c": 50, "justificativa_alternativa_d": 50,
        "justificativa_alternativa_e": 50,
        "nome_professor_comentario": 22, "nome_professor_video": 22,
        "grande_area_oficial": 20, "especialidade": 20, "assunto": 20,
    }
    for i, col in enumerate(COLUNAS, start=1):
        ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = col_widths.get(col, 15)

    # Dados
    for q in questoes:
        row_values = [q.get(col, "") for col in COLUNAS]
        ws.append(row_values)
        row_idx = ws.max_row
        for cell in ws[row_idx]:
            cell.alignment = wrap_align
        ws.row_dimensions[row_idx].height = 60

    # Congela cabeçalho
    ws.freeze_panes = "A2"

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()
